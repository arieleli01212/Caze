"""Creates the 3-slide submission presentation for the HIT Campus Navigation project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE  = RGBColor(0x1A, 0x56, 0x9E)
TEAL  = RGBColor(0x00, 0x7A, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF4, 0xF6, 0xF8)
DARK  = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT= RGBColor(0xFF, 0x6B, 0x35)

W = Inches(13.33)
H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text_box(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def bullet_box(slide, left, top, width, height, items, size=16, color=DARK, indent=Inches(0.3)):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.25), BLUE)
    text_box(slide, Inches(0.4), Inches(0.08), Inches(12), Inches(0.7),
             title, size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, Inches(0.4), Inches(0.8), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)


def footer(slide, label):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), BLUE)
    text_box(slide, Inches(0.4), H - Inches(0.35), Inches(8), Inches(0.32),
             "HIT Campus Navigation  |  Advanced Programming in Java",
             size=10, color=WHITE)
    text_box(slide, Inches(11), H - Inches(0.35), Inches(2), Inches(0.32),
             label, size=10, color=WHITE, align=PP_ALIGN.RIGHT)


# ── slide 1: Functional Description ──────────────────────────────────────────

def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LGRAY

    header_bar(slide,
               "Slide 1 — Functional Description",
               "What the system does, from the user's perspective")
    footer(slide, "1 / 3")

    # Left panel
    add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.55), WHITE)
    text_box(slide, Inches(0.5), Inches(1.45), Inches(5.4), Inches(0.5),
             "Core Features", size=18, bold=True, color=BLUE)

    features = [
        "Find shortest / fewest-stop route between any two HIT campus buildings",
        "Two algorithms selectable by the user:",
        "   Dijkstra  →  FASTEST (weighted path by walking distance)",
        "   BFS  →  FEWEST_SEGMENTS (minimum number of walkways)",
        "Isometric 3D campus view with animated walking pin",
        "Real OpenStreetMap tile view via Leaflet.js (JavaFX WebView)",
        "Route history: browse past queries, clear history",
        "Both views update simultaneously on every route result",
    ]
    bullet_box(slide, Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.7),
               features, size=13)

    # Right panel
    add_rect(slide, Inches(6.4), Inches(1.4), Inches(6.6), Inches(5.55), WHITE)
    text_box(slide, Inches(6.6), Inches(1.45), Inches(6.2), Inches(0.5),
             "User Workflow", size=18, bold=True, color=TEAL)

    steps = [
        "1. Launch server  →  mvn exec:java@server -pl server",
        "2. Launch client  →  mvn javafx:run -pl client",
        "3. Select 'From' and 'To' buildings from drop-downs",
        "4. Choose mode: Fastest (Dijkstra) or Fewest Stops (BFS)",
        "5. Click  Find Route  — route highlights in both views",
        "6. Walker animation traces the path on the 3D canvas",
        "7. Switch to Real Map tab to see the GPS-accurate overlay",
        "8. Open History tab to review or clear past queries",
    ]
    bullet_box(slide, Inches(6.6), Inches(2.0), Inches(6.1), Inches(4.5),
               steps, size=13, color=DARK)

    text_box(slide, Inches(6.6), Inches(5.8), Inches(6.1), Inches(0.9),
             "Campus: 13 HIT buildings, real GPS coordinates, weighted walkways",
             size=12, color=RGBColor(0x55, 0x55, 0x77))


# ── slide 2: HL Architecture ─────────────────────────────────────────────────

def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LGRAY

    header_bar(slide,
               "Slide 2 — HL Architecture",
               "Component overview and design patterns applied")
    footer(slide, "2 / 3")

    # Three columns: algorithm | server | client
    col_w = Inches(3.9)
    gap   = Inches(0.25)
    top   = Inches(1.4)
    ht    = Inches(5.55)

    colors = [RGBColor(0xE8, 0xF4, 0xFD), RGBColor(0xE8, 0xF8, 0xF0), RGBColor(0xFF, 0xF3, 0xE0)]
    borders = [BLUE, TEAL, ACCENT]

    for i, (col_color, border) in enumerate(zip(colors, borders)):
        x = Inches(0.3) + i * (col_w + gap)
        box = add_rect(slide, x, top, col_w, ht, col_color, border)

    labels = ["algorithm  (JAR)", "server", "client  (JavaFX)"]
    label_colors = [BLUE, TEAL, ACCENT]
    contents = [
        [
            "IAlgoShortestPath  (Strategy interface)",
            "AbstractAlgoShortestPath  (Template Method)",
            "DijkstraAlgoShortestPathImpl  — weighted",
            "BFSAlgoShortestPathImpl  — unweighted",
            "Graph  (adjacency list)",
            "",
            "Standalone JAR → added as library",
            "dependency by the server module",
        ],
        [
            "ServerMain  →  Server (Runnable)",
            "ClientHandler  per TCP connection",
            "RequestDispatcher → ControllerFactory",
            "  (Factory Pattern)",
            "",
            "Controllers (IController interface):",
            "  RouteController",
            "  HistoryGetController",
            "  HistoryClearController",
            "",
            "Services:",
            "  NavigationService",
            "  HistoryService",
            "",
            "DAO (IDao interface):",
            "  FileCampusDAO",
            "  FileHistoryDAO",
            "",
            "Protocol: ServerRequest / ServerResponse",
            "  {headers:{action:...}, body:{...}}",
            "Decorator: Scanner + PrintWriter wrappers",
        ],
        [
            "ClientApp  (JavaFX Application)",
            "ClientViewController  (View — MVC)",
            "ClientController  (Controller — MVC)",
            "NavigationClient  (Socket wrapper)",
            "",
            "Views:",
            "  SchematicCanvas  (Canvas + AnimationTimer)",
            "  LeafletMapView  (WebView + Leaflet.js)",
            "  History TableView",
            "",
            "Loosely coupled from server:",
            "  only shared protocol/domain classes",
            "  no direct service or DAO references",
            "",
            "MVC: view fires events to controller;",
            "  controller calls NavigationClient;",
            "  callbacks update UI via Platform.runLater",
        ],
    ]

    for i, (label, items, lc) in enumerate(zip(labels, contents, label_colors)):
        x = Inches(0.3) + i * (col_w + gap)
        text_box(slide, x + Inches(0.15), top + Inches(0.1),
                 col_w - Inches(0.3), Inches(0.42),
                 label, size=15, bold=True, color=lc)
        bullet_box(slide, x + Inches(0.1), top + Inches(0.55),
                   col_w - Inches(0.2), ht - Inches(0.6),
                   items, size=11, color=DARK)

    # Arrow annotations between columns
    def arrow_label(x, text):
        text_box(slide, x, Inches(4.0), Inches(0.95), Inches(0.5),
                 text, size=10, bold=True, color=RGBColor(0x66, 0x66, 0x66),
                 align=PP_ALIGN.CENTER)

    arrow_label(Inches(4.15), "depends\non JAR")
    arrow_label(Inches(8.05), "TCP JSON\nprotocol")


# ── slide 3: Challenges ───────────────────────────────────────────────────────

def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LGRAY

    header_bar(slide,
               "Slide 3 — Challenges",
               "Technical problems encountered and how they were solved")
    footer(slide, "3 / 3")

    challenges = [
        (
            "1  JavaFX Map Rendering",
            "Replacing Swing JXMapViewer2 with a JavaFX-native solution. "
            "Solution: embedded Leaflet.js inside a JavaFX WebView. Building "
            "GPS coordinates are injected from Java into the HTML at page-load "
            "time; route overlays are driven via engine.executeScript() calls.",
        ),
        (
            "2  Module System & JavaFX",
            "JavaFX requires explicit module declarations but the project avoids "
            "module-info.java to keep things simple. Solution: javafx-maven-plugin "
            "0.0.8, which automatically adds --module-path and --add-modules to "
            "the JVM launch command.",
        ),
        (
            "3  Multi-Module Maven Dependency",
            "The algorithm module must be a standalone JAR consumed by the server. "
            "Solution: Maven multi-module parent (packaging=pom) with three child "
            "modules; server pom.xml declares <dependency> on algorithm artifact. "
            "maven-shade-plugin produces an executable fat-jar for the server.",
        ),
        (
            "4  Isometric Projection Math",
            "Rendering 13 buildings in a convincing isometric 3D view required "
            "pure-Java geometry (no game engine). Solution: IsometricProjection "
            "class converts world (x, y) to screen coordinates with a tilt angle, "
            "enabling a smooth 0°→45° slider from top-down to isometric.",
        ),
        (
            "5  Thread Safety in MVC",
            "NavigationClient I/O runs on background threads; JavaFX requires UI "
            "updates on the FX Application Thread. Solution: all callbacks in "
            "ClientViewController are wrapped in Platform.runLater(), keeping the "
            "controller framework-agnostic (no JavaFX imports in ClientController).",
        ),
        (
            "6  TCP Protocol Envelope",
            "Client and server needed a versioned, extensible message format. "
            "Solution: JSON envelope {headers:{action:...}, body:{...}} mirroring "
            "HTTP. ControllerFactory (Factory Pattern) dispatches by action string, "
            "making it trivial to add new API endpoints without modifying existing code.",
        ),
    ]

    rows = 3
    cols = 2
    cw = Inches(6.2)
    ch = Inches(1.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    x0 = Inches(0.3)
    y0 = Inches(1.45)

    for idx, (title, body) in enumerate(challenges):
        row = idx // cols
        col = idx % cols
        x = x0 + col * (cw + gap_x)
        y = y0 + row * (ch + gap_y)
        add_rect(slide, x, y, cw, ch, WHITE, RGBColor(0xCC, 0xCC, 0xCC))
        text_box(slide, x + Inches(0.12), y + Inches(0.08),
                 cw - Inches(0.24), Inches(0.36),
                 title, size=13, bold=True, color=BLUE)
        text_box(slide, x + Inches(0.12), y + Inches(0.44),
                 cw - Inches(0.24), ch - Inches(0.5),
                 body, size=11, color=DARK)


def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1(prs)
    slide2(prs)
    slide3(prs)

    out = "slides.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
